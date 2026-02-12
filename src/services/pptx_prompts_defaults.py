"""Default prompts for PPTX conversion.

These prompts are used as fallback values when prompts are not configured
in the database for a specific profile.
"""

# System prompt for single slide conversion
DEFAULT_SYSTEM_PROMPT = """Generate Python code with convert_to_pptx(html_str, output_path, assets_dir) function.

CRITICAL: Start code with REQUIRED imports:
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

Tools: Presentation, slide_layouts[6], shapes.add_textbox/picture/chart/shape(), Pt(), RGBColor(), PP_ALIGN, Inches(), MSO_SHAPE
Slide: 10" × 7.5". Bounds: left ≥ 0.5", top ≥ 0.5", left + width ≤ 9.5", top + height ≤ 7.0"

COLOR EXTRACTION:
- Extract ALL colors from CSS/inline styles. Convert hex to RGBColor: #102025 → RGBColor(16, 32, 37)
- Use: RGBColor(int(hex[1:3], 16), int(hex[3:5], 16), int(hex[5:7], 16))
- Apply: slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
- Gradients: use first color. Preserve all badge/border/highlight colors.

METRIC CARDS (.metric-card):
- Create VISUAL BOX first: card_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
- Set white background: card_box.fill.solid(); card_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
- Add left border: border_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.05), height); border_shape.fill.solid(); border_shape.fill.fore_color.rgb = RGBColor(r, g, b) from border-left-color CSS; border_shape.line.fill.background()
- Add 3 text boxes INSIDE with 0.2" padding:
  * Label: left + 0.2", top + 0.2", width - 0.4", height 0.3", font Pt(12-14), color #5D6D71, PP_ALIGN.LEFT
  * Value: left + 0.2", top + 0.5", width - 0.4", height 0.8", font Pt(28-36), color #102025, PP_ALIGN.LEFT, bold
  * Subtext: left + 0.2", top + 1.3", width - 0.4", height 0.3", font Pt(12-14), preserve color from CSS, PP_ALIGN.LEFT
- Extract ALL three elements from EVERY metric-card. All cards in grid must have IDENTICAL dimensions.

HIGHLIGHT BOXES (.highlight-box):
- Create RECTANGLE with background color from CSS
- Add left border: border_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.05), height); border_shape.fill.solid(); border_shape.fill.fore_color.rgb = RGBColor(r, g, b)
- Add text box inside with 0.2-0.3" padding, preserve ALL text including <strong> tags

METRIC GRID LAYOUT (.metric-grid):
- 3 columns: Column 1: left=0.5", width=2.8"; Column 2: left=3.5", width=2.8"; Column 3: left=6.5", width=2.8"
- Gap: 0.2" between columns
- Each card: height=1.8-2.0", top=2.8-3.0" (below subtitle)

TWO-COLUMN LAYOUT (.two-column):
- Left: left=0.5", width=4.0"; Right: left=5.0", width=4.0"; Gap: 0.5"

TABLES (<table>):
- Use: table = slide.shapes.add_table(rows, cols, left, top, width, height)
- Extract ALL headers from <thead><tr><th> and ALL rows from <tbody><tr><td>
- Headers: Set background color from CSS (typically #102025), text color white, font Pt(14), bold, PP_ALIGN.LEFT, padding 0.1"
- Cells: Extract ALL text including badges/spans, preserve colors, font Pt(14-16), PP_ALIGN.LEFT for text, PP_ALIGN.RIGHT for numbers, padding 0.1"
- Badges: Extract badge text and color from <span class="lob-badge">, preserve inline styling
- Borders: Set table.first_row = True, apply borders to all cells using cell.fill and cell.line
- Row height: Auto-calculate based on content, minimum 0.3" per row
- Column widths: Distribute evenly or based on content (e.g., 3 columns: 1.5", 1.2", 1.3")
- Position: Within two-column layout, left column typically left=0.5", width=4.0", top=2.5"

TITLE SLIDES (.title-slide class):
- Title: left=Inches(1.0), top=Inches(2.75), width=Inches(8.0), height=Inches(1.2), alignment=PP_ALIGN.CENTER, font=Pt(36), color=RGBColor(255, 255, 255), word_wrap=True
- Subtitle: Calculate top = title_box.top + title_box.height + Inches(0.25), left=Inches(1.0), width=Inches(8.0), height=Inches(0.7), alignment=PP_ALIGN.CENTER, font=Pt(16), color=RGBColor(249, 250, 251), word_wrap=True
- Background: If dark, set slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
- CRITICAL: Title height 1.2" allows 2 lines (36pt ≈ 0.5" per line). Subtitle height 0.7" allows 2 lines (16pt ≈ 0.35" per line). Gap 0.25" prevents overlap.

REGULAR SLIDES:
- Title: top=Inches(0.5), left=Inches(0.5), width=Inches(9.0), height=Inches(1.0), font=Pt(36), PP_ALIGN.LEFT, word_wrap=True
- Subtitle: Calculate top = title_box.top + title_box.height + Inches(0.2), left=Inches(0.5), width=Inches(9.0), height=Inches(0.6), font=Pt(16), PP_ALIGN.LEFT, word_wrap=True
- Body: top = subtitle_box.top + subtitle_box.height + Inches(0.2), left ≥ Inches(0.5), font ≤ Pt(18)
- CRITICAL: Title height 1.0" allows 2 lines (36pt ≈ 0.5" per line). Subtitle height 0.6" allows 2 lines (16pt ≈ 0.3" per line). Gap 0.2" prevents overlap. Subtitle position = title_box.top + title_box.height + gap.

FONT SIZES (CRITICAL):
- ALL titles: EXACTLY Pt(36) - NO exceptions
- ALL subtitles: EXACTLY Pt(16) - NO exceptions
- Set explicitly: title_frame.paragraphs[0].font.size = Pt(36) or Pt(16)

CHARTS:
- CRITICAL: Check for chart images in assets_dir. Chart image filenames may be like "chart_0.png", "chart_overallTrendChart.png", etc.
- If chart images exist (check user prompt for exact filenames): Use slide.shapes.add_picture(os.path.join(assets_dir, filename), left, top, width, height) for EACH image
- Match image filename to canvas element in HTML by canvas ID or position
- Position images at canvas location in HTML. Typical position: left=1.0", top=3.5", width=8.0", height=3.5" (adjust based on HTML layout)
- If no images: Extract Chart.js data from <script> tags.

CONTENT IMAGES (<img> tags):
- HTML may contain <img src="content_image_0.png"> tags — these are real image files in assets_dir
- For EACH: slide.shapes.add_picture(os.path.join(assets_dir, filename), left, top, width, height)
- Position based on the image's location in the HTML layout
- Preserve aspect ratio; typical size: width=2-3", height=1.5-2.5"

Return ONLY Python code."""

# User prompt template for single slide conversion
DEFAULT_USER_PROMPT_TEMPLATE = """Convert this HTML to PowerPoint:

{html_content}

{screenshot_note}

Return Python code with convert_to_pptx(html_str, output_path, assets_dir)."""

# System prompt for multi-slide conversion
DEFAULT_MULTI_SLIDE_SYSTEM_PROMPT = """Generate Python code with add_slide_to_presentation(prs, html_str, assets_dir) function.

CRITICAL: Start code with REQUIRED imports:
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

Tools: prs.slides.add_slide(prs.slide_layouts[6]), shapes.add_textbox/picture/chart/shape(), Pt(), RGBColor(), PP_ALIGN, Inches(), MSO_SHAPE
Slide: 10" × 7.5". Bounds: left ≥ 0.5", top ≥ 0.5", left + width ≤ 9.5", top + height ≤ 7.0"

COLOR EXTRACTION:
- Extract ALL colors from CSS/inline styles. Convert hex to RGBColor: #102025 → RGBColor(16, 32, 37)
- Use: RGBColor(int(hex[1:3], 16), int(hex[3:5], 16), int(hex[5:7], 16))
- Apply: slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
- Gradients: use first color. Preserve all badge/border/highlight colors.

METRIC CARDS (.metric-card):
- Create VISUAL BOX first: card_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
- Set white background: card_box.fill.solid(); card_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
- Add left border: border_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.05), height); border_shape.fill.solid(); border_shape.fill.fore_color.rgb = RGBColor(r, g, b) from border-left-color CSS; border_shape.line.fill.background()
- Add 3 text boxes INSIDE with 0.2" padding:
  * Label: left + 0.2", top + 0.2", width - 0.4", height 0.3", font Pt(12-14), color #5D6D71, PP_ALIGN.LEFT
  * Value: left + 0.2", top + 0.5", width - 0.4", height 0.8", font Pt(28-36), color #102025, PP_ALIGN.LEFT, bold
  * Subtext: left + 0.2", top + 1.3", width - 0.4", height 0.3", font Pt(12-14), preserve color from CSS, PP_ALIGN.LEFT
- Extract ALL three elements from EVERY metric-card. All cards in grid must have IDENTICAL dimensions.

HIGHLIGHT BOXES (.highlight-box):
- Create RECTANGLE with background color from CSS
- Add left border: border_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.05), height); border_shape.fill.solid(); border_shape.fill.fore_color.rgb = RGBColor(r, g, b)
- Add text box inside with 0.2-0.3" padding, preserve ALL text including <strong> tags

METRIC GRID LAYOUT (.metric-grid):
- 3 columns: Column 1: left=0.5", width=2.8"; Column 2: left=3.5", width=2.8"; Column 3: left=6.5", width=2.8"
- Gap: 0.2" between columns
- Each card: height=1.8-2.0", top=2.8-3.0" (below subtitle)

TWO-COLUMN LAYOUT (.two-column):
- Left: left=0.5", width=4.0"; Right: left=5.0", width=4.0"; Gap: 0.5"

TABLES (<table>):
- Use: table = slide.shapes.add_table(rows, cols, left, top, width, height)
- Extract ALL headers from <thead><tr><th> and ALL rows from <tbody><tr><td>
- Headers: Set background color from CSS (typically #102025), text color white, font Pt(14), bold, PP_ALIGN.LEFT, padding 0.1"
- Cells: Extract ALL text including badges/spans, preserve colors, font Pt(14-16), PP_ALIGN.LEFT for text, PP_ALIGN.RIGHT for numbers, padding 0.1"
- Badges: Extract badge text and color from <span class="lob-badge">, preserve inline styling
- Borders: Set table.first_row = True, apply borders to all cells using cell.fill and cell.line
- Row height: Auto-calculate based on content, minimum 0.3" per row
- Column widths: Distribute evenly or based on content (e.g., 3 columns: 1.5", 1.2", 1.3")
- Position: Within two-column layout, left column typically left=0.5", width=4.0", top=2.5"

TITLE SLIDES (.title-slide class):
- Title: left=Inches(1.0), top=Inches(2.75), width=Inches(8.0), height=Inches(1.2), alignment=PP_ALIGN.CENTER, font=Pt(36), color=RGBColor(255, 255, 255), word_wrap=True
- Subtitle: Calculate top = title_box.top + title_box.height + Inches(0.25), left=Inches(1.0), width=Inches(8.0), height=Inches(0.7), alignment=PP_ALIGN.CENTER, font=Pt(16), color=RGBColor(249, 250, 251), word_wrap=True
- Background: If dark, set slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
- CRITICAL: Title height 1.2" allows 2 lines (36pt ≈ 0.5" per line). Subtitle height 0.7" allows 2 lines (16pt ≈ 0.35" per line). Gap 0.25" prevents overlap.

REGULAR SLIDES:
- Title: top=Inches(0.5), left=Inches(0.5), width=Inches(9.0), height=Inches(1.0), font=Pt(36), PP_ALIGN.LEFT, word_wrap=True
- Subtitle: Calculate top = title_box.top + title_box.height + Inches(0.2), left=Inches(0.5), width=Inches(9.0), height=Inches(0.6), font=Pt(16), PP_ALIGN.LEFT, word_wrap=True
- Body: top = subtitle_box.top + subtitle_box.height + Inches(0.2), left ≥ Inches(0.5), font ≤ Pt(18)
- CRITICAL: Title height 1.0" allows 2 lines (36pt ≈ 0.5" per line). Subtitle height 0.6" allows 2 lines (16pt ≈ 0.3" per line). Gap 0.2" prevents overlap. Subtitle position = title_box.top + title_box.height + gap.

FONT SIZES (CRITICAL):
- ALL titles: EXACTLY Pt(36) - NO exceptions
- ALL subtitles: EXACTLY Pt(16) - NO exceptions
- Set explicitly: title_frame.paragraphs[0].font.size = Pt(36) or Pt(16)

CHARTS:
- CRITICAL: Check for chart images in assets_dir. Chart image filenames may be like "chart_0.png", "chart_overallTrendChart.png", etc.
- If chart images exist (check user prompt for exact filenames): Use slide.shapes.add_picture(os.path.join(assets_dir, filename), left, top, width, height) for EACH image
- Match image filename to canvas element in HTML by canvas ID or position
- Position images at canvas location in HTML. Typical position: left=1.0", top=3.5", width=8.0", height=3.5" (adjust based on HTML layout)
- If no images: Extract Chart.js data from <script> tags.

CONTENT IMAGES (<img> tags):
- HTML may contain <img src="content_image_0.png"> tags — these are real image files in assets_dir
- For EACH: slide.shapes.add_picture(os.path.join(assets_dir, filename), left, top, width, height)
- Position based on the image's location in the HTML layout
- Preserve aspect ratio; typical size: width=2-3", height=1.5-2.5"

Return ONLY Python code."""

# User prompt template for multi-slide conversion
DEFAULT_MULTI_SLIDE_USER_PROMPT = """Add slide from HTML to presentation:

{html_content}

{screenshot_note}

Return Python code with add_slide_to_presentation(prs, html_str, assets_dir)."""
