# src/services/converter_jail/pptx_runner.py
"""In-jail trusted runner for PPTX conversion (SDR-4437 HIGH-5).

Runs the whole per-slide execution loop inside the subprocess jail. The live
python-pptx Presentation never crosses the process boundary. Per-slide fault
isolation and per-slide progress reporting (one stdout line per slide) are
preserved. Generated code's stdout is redirected to stderr so it cannot spoof
the trusted progress channel.

CLI: python -I pptx_runner.py <src_parent> <job_dir> <output_path>
(<src_parent> is part of the launcher's CLI contract; this runner does NOT
put it on sys.path — importing anything via `src.services` would execute that
package's heavy __init__ and drag the credential-capable SDK graph into the
jail child. Siblings are imported as top-level modules instead.)
"""

import importlib.util
import json
import os
import sys
import tempfile
from pathlib import Path
from typing import Callable, Optional

_HERE = Path(__file__).resolve().parent

if __package__ in (None, ""):
    # Script mode (inside the jail child): import siblings as top-level
    # modules so src/services/__init__.py (SDK/agent graph) never executes.
    sys.path.insert(0, str(_HERE))
    import protocol
    from codeprep import sanitize_code
else:
    # Imported on the trusted host (unit tests): normal package imports.
    from src.services.converter_jail import protocol
    from src.services.converter_jail.codeprep import sanitize_code


_REQUIRED_IMPORTS = """from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
"""


def _add_fallback(prs, slide_number: int, title: str) -> None:
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.color.rgb = RGBColor(16, 32, 37)


def _exec_snippet(code: str, prs, html_str: str, assets_dir: str) -> bool:
    """Execute one generated snippet against prs. Returns True on success."""
    code = sanitize_code(code)
    if "from pptx.util import" not in code:
        code = _REQUIRED_IMPORTS + "\n" + code

    tmp = Path(tempfile.mktemp(suffix=".py", prefix="jail_slide_"))
    tmp.write_text(code, encoding="utf-8")
    try:
        spec = importlib.util.spec_from_file_location("jail_slide", str(tmp))
        module = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(module)
        module.add_slide_to_presentation(prs, html_str, assets_dir)
        return True
    except Exception:
        return False
    finally:
        if tmp.exists():
            tmp.unlink()


def run(job_dir: str, output_path: str, emit: Optional[Callable[[str], None]] = None) -> None:
    """Build the deck. `emit` receives each progress LINE (incl. newline);
    defaults to writing to the real stdout fd. Importable for unit tests."""
    from pptx import Presentation
    from pptx.util import Inches

    if emit is None:
        # Duplicate the real stdout, then redirect Python-level stdout to stderr
        # so generated print()s cannot appear on the progress channel.
        real_fd = os.dup(1)
        os.dup2(2, 1)
        sys.stdout = sys.stderr

        def emit(line: str) -> None:
            os.write(real_fd, line.encode("utf-8"))

    manifest = json.loads((Path(job_dir) / protocol.MANIFEST_NAME).read_text())
    slides = manifest["slides"]
    total = len(slides)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for entry in slides:
        i = entry["index"] + 1
        sdir = Path(job_dir) / entry["dir"]
        assets_dir = str(sdir / protocol.ASSETS_DIR)

        if not entry.get("has_code"):
            _add_fallback(prs, i, f"Slide {i}")
        else:
            html_str = (sdir / protocol.HTML_NAME).read_text(encoding="utf-8")
            code = (sdir / protocol.CODE_NAME).read_text(encoding="utf-8")
            if not _exec_snippet(code, prs, html_str, assets_dir):
                _add_fallback(prs, i, "Slide Content")

        emit(protocol.encode_progress(i, total, f"Building slide {i}/{total}…"))

    prs.save(output_path)


if __name__ == "__main__":
    # argv: [runner, src_parent, job_dir, output_path]
    run(sys.argv[2], sys.argv[3])
