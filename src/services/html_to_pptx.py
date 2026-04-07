#!/usr/bin/env python3

import asyncio
import base64
import importlib.util
import logging
import os
import re
import tempfile
import time
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple

from bs4 import BeautifulSoup
from databricks.sdk import WorkspaceClient
from pptx import Presentation
from pptx.util import Inches

from src.services.pptx_prompts_defaults import (
    DEFAULT_SYSTEM_PROMPT,
    DEFAULT_USER_PROMPT_TEMPLATE,
    DEFAULT_MULTI_SLIDE_SYSTEM_PROMPT,
    DEFAULT_MULTI_SLIDE_USER_PROMPT,
)

logger = logging.getLogger(__name__)

# Path to the placeholder PPTX whose first slide is always prepended.
# Lives under src/assets/ so it's included in Databricks Apps deployments.
PLACEHOLDER_PPTX_PATH = Path(__file__).resolve().parent.parent / "assets" / "placeholder.pptx"

MAX_CONCURRENT_LLM = 5


class PPTXConversionError(Exception):
    """Raised when PPTX conversion fails."""
    pass


class HtmlToPptxConverterV3:
    """V3 converter using maximum LLM flexibility approach.
    
    This converter uses an LLM to analyze HTML and generate Python code
    that creates PowerPoint slides. This approach provides maximum flexibility
    and handles diverse layouts automatically.
    
    Attributes:
        model_endpoint: LLM model endpoint name
        ws_client: Databricks workspace client
        llm_client: OpenAI-compatible client for LLM calls
    """
    
    # Model configuration
    DEFAULT_MODEL = "databricks-claude-sonnet-4-5"
    

    def __init__(
        self,
        workspace_client: Optional[WorkspaceClient] = None,
        model_endpoint: Optional[str] = None,
    ):
        """Initialize V3 converter.
        
        Args:
            workspace_client: Databricks client (optional, uses singleton if not provided)
            model_endpoint: LLM model name (default: databricks-claude-sonnet-4-5)
        """
        self.model_endpoint = model_endpoint or self.DEFAULT_MODEL
        
        # Use provided client or get singleton from databricks_client
        if workspace_client:
            self.ws_client = workspace_client
        else:
            from src.core.databricks_client import get_databricks_client
            self.ws_client = get_databricks_client()
        
        # Initialize OpenAI-compatible client for LLM calls
        self.llm_client = self.ws_client.serving_endpoints.get_open_ai_client()
        
        # Load default prompts
        self.SYSTEM_PROMPT = DEFAULT_SYSTEM_PROMPT
        self.USER_PROMPT_TEMPLATE = DEFAULT_USER_PROMPT_TEMPLATE
        self.MULTI_SLIDE_SYSTEM_PROMPT = DEFAULT_MULTI_SLIDE_SYSTEM_PROMPT
        self.MULTI_SLIDE_USER_PROMPT = DEFAULT_MULTI_SLIDE_USER_PROMPT
        
        logger.info(
            "V3 Converter initialized",
            extra={"model": self.model_endpoint}
        )
    
    async def convert_html_to_pptx(
        self,
        html_str: str,
        output_path: str,
        use_screenshot: bool = True,
        html_source_path: Optional[str] = None
    ) -> str:
        """Convert single HTML slide to PowerPoint.
        
        Args:
            html_str: HTML content
            output_path: Path to save PPTX
            use_screenshot: Whether to capture and use screenshot
            html_source_path: Path to HTML file (for screenshot)
        
        Returns:
            Path to created PPTX file
        
        Raises:
            PPTXConversionError: If conversion fails
        """
        logger.info("Converting single HTML to PowerPoint")
        
        # 1. Setup working directory
        work_dir = Path(tempfile.mkdtemp(prefix="v3_convert_"))
        assets_dir = work_dir / "assets"
        assets_dir.mkdir()
        
        # 2. Chart screenshots are now captured client-side
        # Server-side Playwright capture has been removed
        chart_images = []
        if use_screenshot and html_source_path:
            logger.warning(
                "Server-side screenshot capture is no longer supported. "
                "Use client-side capture via chart_images parameter."
            )
            print("[PPTX_CONVERTER] WARNING: Server-side screenshots disabled. Use client-side capture.")
        
        # 2b. Extract base64 content images before sending to LLM
        html_str, content_images = self._extract_and_save_content_images(html_str, str(assets_dir))
        chart_images.extend(content_images)

        # 3. Call LLM to generate converter code
        logger.info("Calling LLM to generate converter code")
        converter_code = await self._generate_converter_code(
            html_str,
            chart_images=chart_images
        )
        
        if not converter_code:
            raise PPTXConversionError("Failed to generate converter code from LLM")
        
        # 4. Execute generated code
        logger.info("Executing generated converter")
        self._execute_single_slide_converter(
            converter_code,
            html_str,
            output_path,
            str(assets_dir)
        )
        
        logger.info("PowerPoint created", extra={"path": output_path})
        return output_path
    
    async def convert_slide_deck(
        self,
        slides: List[str],
        output_path: str,
        use_screenshot: bool = True,
        html_source_paths: Optional[List[str]] = None,
        chart_images_per_slide: Optional[List[Dict[str, str]]] = None,
        progress_callback: Optional[Callable[[int, int, str], None]] = None,
    ) -> str:
        """Convert multiple HTML slides to PowerPoint deck.

        Uses a two-phase approach:
          Phase 1 — Prepare assets + parallel LLM code generation
          Phase 2 — Sequential execution against a single Presentation object

        Args:
            slides: List of HTML strings
            output_path: Path to save PPTX
            use_screenshot: Whether to capture screenshots
            html_source_paths: Paths to HTML files (for screenshots, if not using client images)
            chart_images_per_slide: List of dicts mapping canvas_id to base64 data URL per slide
            progress_callback: Optional ``(current, total, status)`` callback.
        
        Returns:
            Path to created PPTX file
        
        Raises:
            PPTXConversionError: If conversion fails
        """
        total = len(slides)
        log_msg = (
            f"Converting {total} slides to PowerPoint - "
            f"total HTML size: {sum(len(html) for html in slides)}, "
            f"use_screenshot: {use_screenshot}"
        )
        logger.info(log_msg)
        print(f"[PPTX_CONVERTER] {log_msg}")

        # -- Phase 1: Prepare assets + parallel LLM code generation ----------
        if progress_callback:
            try:
                progress_callback(0, total, f"Generating code for {total} slides...")
            except Exception:
                pass

        slide_inputs: List[Tuple[str, List[str], List[str], str]] = []
        for i, html_str in enumerate(slides, 1):
            chart_imgs = None
            if chart_images_per_slide and i - 1 < len(chart_images_per_slide):
                chart_imgs = chart_images_per_slide[i - 1]
            prepared = self._prepare_slide(html_str, chart_imgs, i)
            slide_inputs.append(prepared)

        def _on_codegen_progress(completed: int, codegen_total: int) -> None:
            if progress_callback:
                try:
                    progress_callback(
                        completed, codegen_total,
                        f"Generating code: {completed}/{codegen_total} slides ready…",
                    )
                except Exception:
                    pass

        t0 = time.time()
        codes = await self._generate_all_codes(slide_inputs, on_codegen_progress=_on_codegen_progress)
        logger.info(
            "Parallel codegen complete",
            extra={"total_slides": total, "duration_s": f"{time.time() - t0:.1f}"},
        )

        # -- Phase 2: Sequential execution -----------------------------------
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        for i, (code, (html_str, _chart_files, _content_files, assets_dir)) in enumerate(
            zip(codes, slide_inputs), 1,
        ):
            print(f"[PPTX_CONVERTER] Executing slide {i}/{total}")
            self._execute_slide(code, prs, html_str, assets_dir, i)

            if progress_callback:
                try:
                    progress_callback(i, total, f"Building slide {i}/{total}…")
                except Exception:
                    pass

        prs.save(output_path)
        logger.info(
            "PowerPoint deck created",
            extra={"path": output_path, "slide_count": total},
        )
        return output_path
    
    async def _add_slide_to_presentation(
        self,
        prs: Presentation,
        html_str: str,
        use_screenshot: bool,
        html_source_path: Optional[str],
        slide_number: int,
        client_chart_images: Optional[Dict[str, str]] = None
    ) -> None:
        """Add a slide to existing presentation using V3 approach.
        
        Args:
            prs: Presentation object to add slide to
            html_str: HTML content for the slide
            use_screenshot: Whether to capture and use screenshot
            html_source_path: Path to HTML file (for screenshot capture, if not using client images)
            slide_number: Slide number for logging purposes
            client_chart_images: Dict mapping canvas_id to base64 data URL (from client-side capture)
        """
        # 1. Setup working directory for this slide
        work_dir = Path(tempfile.mkdtemp(prefix=f"v3_slide_{slide_number}_"))
        assets_dir = work_dir / "assets"
        assets_dir.mkdir()
        
        # 2. Get chart images - prefer client-provided, fallback to server-side capture
        chart_images = []
        if use_screenshot:
            if client_chart_images:
                # Save client-provided base64 images to files
                logger.info(
                    "Saving client-provided chart images",
                    extra={
                        "slide_number": slide_number,
                        "chart_count": len(client_chart_images),
                        "assets_dir": str(assets_dir)
                    }
                )
                print(f"[PPTX_CONVERTER] Saving {len(client_chart_images)} client-provided chart images for slide {slide_number}")
                chart_images = self._save_client_chart_images(client_chart_images, str(assets_dir))
                if chart_images:
                    logger.info(
                        "Client chart images saved successfully",
                        extra={
                            "slide_number": slide_number,
                            "chart_count": len(chart_images),
                            "chart_files": chart_images
                        }
                    )
                    print(f"[PPTX_CONVERTER] Saved {len(chart_images)} chart images for slide {slide_number}: {chart_images}")
            elif html_source_path:
                # Server-side Playwright capture has been removed
                # Client-side capture should be used instead
                logger.warning(
                    "Server-side screenshot capture is no longer supported for slide",
                    extra={
                        "slide_number": slide_number,
                        "html_source_path": html_source_path
                    }
                )
                print(f"[PPTX_CONVERTER] WARNING: Server-side screenshots disabled for slide {slide_number}. Use client-side capture.")
            else:
                logger.warning("No chart screenshots captured", extra={"slide_number": slide_number})
                print(f"[PPTX_CONVERTER] WARNING: No chart screenshots captured for slide {slide_number}")

        # 2b. Extract base64 content images before sending to LLM
        has_base64_img = "data:image" in html_str
        print(f"[PPTX_CONVERTER] Slide {slide_number}: base64 images in HTML: {has_base64_img}, HTML length: {len(html_str)}")
        html_str, content_images = self._extract_and_save_content_images(html_str, str(assets_dir))
        if content_images:
            print(f"[PPTX_CONVERTER] Slide {slide_number}: extracted {len(content_images)} content images: {content_images}")

        # 3. Call LLM to generate code for adding this slide
        logger.debug("Calling LLM for slide", extra={"slide_number": slide_number})
        slide_code = await self._generate_slide_adder_code(
            html_str,
            chart_images=chart_images,
            content_images=content_images,
        )
        
        if not slide_code:
            logger.warning(
                "Failed to generate code for slide, skipping",
                extra={"slide_number": slide_number}
            )
            return
        
        # 4. Execute generated code to add slide
        logger.debug("Adding slide to presentation", extra={"slide_number": slide_number})
        self._execute_slide_adder(
            slide_code,
            prs,
            html_str,
            str(assets_dir)
        )
        
        logger.debug("Slide added", extra={"slide_number": slide_number})
    
    async def _generate_converter_code(
        self,
        html_str: str,
        chart_images: list[str]
    ) -> Optional[str]:
        """Call LLM to generate converter code for single slide.
        
        Args:
            html_str: HTML content to convert
            chart_images: List of chart image filenames (e.g., ['chart_0.png', 'chart_1.png'])
        
        Returns:
            Generated Python code or None if generation fails
        """
        logger.debug(
            "Generating converter code",
            extra={
                "original_html_length": len(html_str),
                "chart_images": chart_images,
                "chart_count": len(chart_images),
            }
        )
        
        # Truncate HTML if too long
        html_content = self._truncate_html(html_str)
        
        logger.debug(
            "HTML truncated for LLM",
            extra={
                "original_length": len(html_str),
                "truncated_length": len(html_content),
                "was_truncated": len(html_content) < len(html_str),
                "truncated_preview": html_content[:500] + "..." if len(html_content) > 500 else html_content,
            }
        )
        
        screenshot_note = ""
        if chart_images:
            chart_files_str = ", ".join(chart_images)
            # Provide clear instructions with all available chart files
            if len(chart_images) == 1:
                screenshot_note = (
                    f"Chart image available: {chart_files_str}. "
                    f"CRITICAL: Use slide.shapes.add_picture(os.path.join(assets_dir, '{chart_images[0]}'), left, top, width, height) to add this image. "
                    f"Find the canvas element in the HTML (look for <canvas> tags) and position the image at that location. "
                    f"Typical position: left=1.0\", top=3.5\", width=8.0\", height=3.5\" (adjust based on HTML layout)."
                )
            else:
                screenshot_note = (
                    f"Chart images available: {chart_files_str}. "
                    f"CRITICAL: Use slide.shapes.add_picture() for each image. "
                    f"Match each image filename to its corresponding canvas element in the HTML by canvas ID or position. "
                    f"Find canvas elements in HTML and position images at their locations. "
                    f"Example: slide.shapes.add_picture(os.path.join(assets_dir, '{chart_images[0]}'), left, top, width, height)"
                )
        else:
            screenshot_note = "No chart images. Extract Chart.js data from <script> tags. Only use data from HTML."
        
        prompt = self.USER_PROMPT_TEMPLATE.format(
            html_content=html_content,
            screenshot_note=screenshot_note
        )
        
        return await self._call_llm(self.SYSTEM_PROMPT, prompt)
    
    async def _generate_slide_adder_code(
        self,
        html_str: str,
        chart_images: list[str],
        content_images: Optional[list[str]] = None,
    ) -> Optional[str]:
        """Call LLM to generate code for adding slide to existing presentation.

        Args:
            html_str: HTML content to convert
            chart_images: List of chart image filenames (e.g., ['chart_0.png', 'chart_1.png'])
            content_images: List of content image filenames (e.g., ['content_image_0.png'])

        Returns:
            Generated Python code or None if generation fails
        """
        content_images = content_images or []
        logger.debug(
            "Generating slide adder code",
            extra={
                "original_html_length": len(html_str),
                "chart_images": chart_images,
                "chart_count": len(chart_images),
                "content_images": content_images,
                "content_image_count": len(content_images),
            }
        )

        # Truncate HTML if too long
        html_content = self._truncate_html(html_str)

        logger.debug(
            "HTML truncated for LLM (slide adder)",
            extra={
                "original_length": len(html_str),
                "truncated_length": len(html_content),
                "was_truncated": len(html_content) < len(html_str),
                "truncated_preview": html_content[:500] + "..." if len(html_content) > 500 else html_content,
            }
        )

        screenshot_note = ""
        if chart_images:
            chart_files_str = ", ".join(chart_images)
            # Provide clear instructions with all available chart files
            if len(chart_images) == 1:
                screenshot_note = (
                    f"Chart image available: {chart_files_str}. "
                    f"CRITICAL: Use slide.shapes.add_picture(os.path.join(assets_dir, '{chart_images[0]}'), left, top, width, height) to add this image. "
                    f"Find the canvas element in the HTML (look for <canvas> tags) and position the image at that location. "
                    f"Typical position: left=1.0\", top=3.5\", width=8.0\", height=3.5\" (adjust based on HTML layout)."
                )
            else:
                screenshot_note = (
                    f"Chart images available: {chart_files_str}. "
                    f"CRITICAL: Use slide.shapes.add_picture() for each image. "
                    f"Match each image filename to its corresponding canvas element in the HTML by canvas ID or position. "
                    f"Find canvas elements in HTML and position images at their locations. "
                    f"Example: slide.shapes.add_picture(os.path.join(assets_dir, '{chart_images[0]}'), left, top, width, height)"
                )
        else:
            screenshot_note = "No chart images. Extract Chart.js data from <script> tags. Only use data from HTML."

        if content_images:
            content_files_str = ", ".join(content_images)
            screenshot_note += (
                f"\nContent images in assets_dir: {content_files_str}. "
                f"These are logos/icons extracted from <img> tags in the HTML. "
                f"CRITICAL: You MUST add each content image using slide.shapes.add_picture(os.path.join(assets_dir, filename), left, top, width, height). "
                f"Match position and size to the original <img> tag's placement in the HTML (e.g. a logo in the top-right corner). "
                f"Typical logo size: width=Inches(1.0), height=Inches(0.5). Place in the header area near the title."
            )

        print(f"[PPTX_CONVERTER] screenshot_note for LLM: {screenshot_note[:500]}")

        prompt = self.MULTI_SLIDE_USER_PROMPT.format(
            html_content=html_content,
            screenshot_note=screenshot_note
        )

        return await self._call_llm(self.MULTI_SLIDE_SYSTEM_PROMPT, prompt)

    # -- Slide prep / parallel codegen / execution helpers --------------------

    def _prepare_slide(
        self, html_str: str, client_chart_images: Optional[Dict[str, str]], slide_num: int,
    ) -> Tuple[str, List[str], List[str], str]:
        """Save images to disk and clean HTML. Returns (html, chart_files, content_files, assets_dir)."""
        work_dir = Path(tempfile.mkdtemp(prefix=f"v3_slide_{slide_num}_"))
        assets_dir = work_dir / "assets"
        assets_dir.mkdir()

        chart_images: List[str] = []
        if client_chart_images:
            chart_images = self._save_client_chart_images(client_chart_images, str(assets_dir))
            if chart_images:
                print(f"[PPTX_CONVERTER] Slide {slide_num}: saved {len(chart_images)} chart images")

        has_base64_img = "data:image" in html_str
        print(f"[PPTX_CONVERTER] Slide {slide_num}: base64 images in HTML: {has_base64_img}, HTML length: {len(html_str)}")
        html_str, content_images = self._extract_and_save_content_images(html_str, str(assets_dir))
        if content_images:
            print(f"[PPTX_CONVERTER] Slide {slide_num}: extracted {len(content_images)} content images: {content_images}")
        print(f"[PPTX_CONVERTER] Slide {slide_num}: HTML length after image extraction: {len(html_str)}")

        return html_str, chart_images, content_images, str(assets_dir)

    def _generate_code_sync(
        self, html_str: str, chart_images: List[str], content_images: Optional[List[str]] = None,
    ) -> Optional[str]:
        """Synchronous code generation for use with asyncio.to_thread."""
        content_images = content_images or []
        html_content = self._truncate_html(html_str)

        screenshot_note = self._build_screenshot_note(chart_images, content_images)

        prompt = self.MULTI_SLIDE_USER_PROMPT.format(
            html_content=html_content,
            screenshot_note=screenshot_note,
        )
        return self._call_llm_sync(self.MULTI_SLIDE_SYSTEM_PROMPT, prompt)

    def _build_screenshot_note(
        self, chart_images: List[str], content_images: Optional[List[str]] = None,
    ) -> str:
        """Build the screenshot/asset note portion of the user prompt."""
        content_images = content_images or []
        screenshot_note = ""
        if chart_images:
            chart_files_str = ", ".join(chart_images)
            if len(chart_images) == 1:
                screenshot_note = (
                    f"Chart image available: {chart_files_str}. "
                    f"CRITICAL: Use slide.shapes.add_picture(os.path.join(assets_dir, '{chart_images[0]}'), left, top, width, height) to add this image. "
                    f"Find the canvas element in the HTML (look for <canvas> tags) and position the image at that location. "
                    f"Typical position: left=1.0\", top=3.5\", width=8.0\", height=3.5\" (adjust based on HTML layout)."
                )
            else:
                screenshot_note = (
                    f"Chart images available: {chart_files_str}. "
                    f"CRITICAL: Use slide.shapes.add_picture() for each image. "
                    f"Match each image filename to its corresponding canvas element in the HTML by canvas ID or position. "
                    f"Find canvas elements in HTML and position images at their locations. "
                    f"Example: slide.shapes.add_picture(os.path.join(assets_dir, '{chart_images[0]}'), left, top, width, height)"
                )
        else:
            screenshot_note = "No chart images. Extract Chart.js data from <script> tags. Only use data from HTML."

        if content_images:
            content_files_str = ", ".join(content_images)
            screenshot_note += (
                f"\nContent images in assets_dir: {content_files_str}. "
                f"These are logos/icons extracted from <img> tags in the HTML. "
                f"CRITICAL: You MUST add each content image using slide.shapes.add_picture(os.path.join(assets_dir, filename), left, top, width, height). "
                f"Match position and size to the original <img> tag's placement in the HTML (e.g. a logo in the top-right corner). "
                f"Typical logo size: width=Inches(1.0), height=Inches(0.5). Place in the header area near the title."
            )
        return screenshot_note

    async def _generate_all_codes(
        self,
        slide_inputs: List[Tuple[str, List[str], List[str], str]],
        on_codegen_progress: Optional[Callable[[int, int], None]] = None,
    ) -> List[Optional[str]]:
        """Generate code for all slides in parallel, capped at MAX_CONCURRENT_LLM.

        Args:
            slide_inputs: List of (html, chart_files, content_files, assets_dir).
            on_codegen_progress: Optional ``(completed, total)`` callback fired
                after each slide's code is generated.
        """
        sem = asyncio.Semaphore(MAX_CONCURRENT_LLM)
        total = len(slide_inputs)
        completed = 0
        lock = asyncio.Lock()

        async def gen_one(html_str: str, chart_images: List[str], content_images: List[str]) -> Optional[str]:
            nonlocal completed
            async with sem:
                code = await asyncio.to_thread(
                    self._generate_code_sync, html_str, chart_images, content_images,
                )
            async with lock:
                completed += 1
                if on_codegen_progress:
                    try:
                        on_codegen_progress(completed, total)
                    except Exception:
                        pass
            return code

        return list(await asyncio.gather(
            *(gen_one(h, c, ci) for h, c, ci, _ad in slide_inputs)
        ))

    def _execute_slide(
        self, code: Optional[str], prs: Presentation, html_str: str,
        assets_dir: str, slide_number: int,
    ) -> None:
        """Execute generated code for one slide, with fallback on failure."""
        if not code:
            logger.warning("No code for slide %d, adding fallback", slide_number)
            self._add_fallback_slide(prs, slide_number)
            return
        self._execute_slide_adder(code, prs, html_str, assets_dir)

    def _add_fallback_slide(self, prs: Presentation, slide_number: int) -> None:
        """Add a placeholder fallback slide when code generation or execution fails."""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = f"Slide {slide_number}"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.color.rgb = RGBColor(16, 32, 37)

    def _call_llm_sync(self, system_prompt: str, user_prompt: str) -> Optional[str]:
        """Synchronous LLM call — core implementation used by both async and threaded paths."""
        start = time.time()
        try:
            response = self.llm_client.chat.completions.create(
                model=self.model_endpoint,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.2,
                max_tokens=16384,
                timeout=300,
                extra_body={
                    "thinking": {
                        "type": "enabled",
                        "budget_tokens": 10240,
                    }
                },
            )
            duration = time.time() - start
            logger.info("LLM call completed", extra={"duration_s": f"{duration:.1f}"})

            code = self._extract_text_content(response.choices[0].message.content)
            return self._strip_markdown_fences(code) if code else None

        except Exception as e:
            duration = time.time() - start
            logger.error("LLM call failed", extra={"duration_s": f"{duration:.1f}", "error": str(e)}, exc_info=True)
            return None

    async def _call_llm(self, system_prompt: str, user_prompt: str) -> Optional[str]:
        """Async wrapper — delegates to _call_llm_sync."""
        return self._call_llm_sync(system_prompt, user_prompt)
    
    @staticmethod
    def _extract_text_content(content) -> str:
        """Extract text from LLM response, handling reasoning model responses.

        When reasoning/thinking is enabled, ``content`` is a list of blocks
        (e.g. [{"type": "reasoning", ...}, {"type": "text", "text": "..."}]).
        Otherwise it's a plain string.
        """
        if isinstance(content, str):
            return content
        if isinstance(content, list):
            for block in content:
                if isinstance(block, dict) and block.get("type") == "text":
                    return block.get("text", "")
        # Fallback — stringify whatever we got
        return str(content) if content else ""

    @staticmethod
    def _strip_markdown_fences(code: str) -> str:
        """Remove markdown code fences from LLM-generated code.

        Handles variations like ```python, ```Python, ``` python, bare ```,
        trailing ```, and leading/trailing whitespace around fences.
        """
        # Try to extract content between ```python ... ``` (case-insensitive, optional whitespace)
        match = re.search(
            r"```[Pp]ython\s*\n(.*?)```", code, re.DOTALL
        )
        if match:
            return match.group(1).strip()

        # Try bare ``` ... ```
        match = re.search(r"```\s*\n(.*?)```", code, re.DOTALL)
        if match:
            return match.group(1).strip()

        # If the code just starts/ends with ``` on its own line, strip them
        lines = code.strip().splitlines()
        if lines and lines[0].strip().startswith("```"):
            lines = lines[1:]
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        return "\n".join(lines)

    def _save_client_chart_images(
        self, 
        client_chart_images: Dict[str, str], 
        assets_dir: str
    ) -> list[str]:
        """Save client-provided base64 chart images to files.
        
        Args:
            client_chart_images: Dict mapping canvas_id to base64 data URL
            assets_dir: Directory to save chart images
        
        Returns:
            List of filenames of saved chart images (e.g., ['chart_0.png', 'chart_1.png'])
        """
        import base64
        
        chart_images = []
        assets_path = Path(assets_dir)
        
        for canvas_id, base64_data in client_chart_images.items():
            try:
                # Extract base64 data (remove data:image/png;base64, prefix if present)
                if ',' in base64_data:
                    base64_data = base64_data.split(',', 1)[1]
                
                # Decode base64 to bytes
                image_bytes = base64.b64decode(base64_data)
                
                # Determine filename from canvas_id
                # Canvas IDs might be like "chart_0", "my_chart", etc.
                if canvas_id.startswith('chart_'):
                    filename = f"{canvas_id}.png"
                else:
                    # Use canvas_id as filename, sanitize if needed
                    safe_id = "".join(c if c.isalnum() or c in ('_', '-') else '_' for c in canvas_id)
                    filename = f"chart_{safe_id}.png"
                
                chart_path = assets_path / filename
                chart_path.write_bytes(image_bytes)
                
                if chart_path.exists():
                    file_size = chart_path.stat().st_size
                    logger.info(
                        f"Saved client chart image: {filename}",
                        extra={
                            "canvas_id": canvas_id,
                            "filename": filename,
                            "file_size": file_size
                        }
                    )
                    print(f"[CLIENT_CHART] Saved chart image: {filename} ({file_size} bytes) from canvas_id: {canvas_id}")
                    chart_images.append(filename)
                else:
                    logger.warning(f"Chart image file not created: {filename}")
                    print(f"[CLIENT_CHART] WARNING: Chart image file not created: {filename}")
            except Exception as e:
                logger.warning(f"Failed to save client chart image {canvas_id}: {e}", exc_info=True)
                print(f"[CLIENT_CHART] WARNING: Failed to save chart image for canvas_id {canvas_id}: {e}")
        
        return chart_images
    
    # Regex matching <img> tags whose src is a base64 data URI.
    _BASE64_IMG_RE = re.compile(
        r'(<img\b[^>]*?\bsrc=")data:image/(png|jpeg|jpg|gif|svg\+xml);base64,([A-Za-z0-9+/=\s]+)(")',
        re.IGNORECASE,
    )

    # Map MIME sub-type to file extension
    _EXT_MAP = {"png": ".png", "jpeg": ".jpg", "jpg": ".jpg", "gif": ".gif", "svg+xml": ".svg"}

    @staticmethod
    def _svg_to_png(svg_bytes: bytes) -> bytes:
        """Convert SVG bytes to PNG using svgpathtools + Pillow (pure Python).

        Parses SVG ``<path>`` elements, samples bezier curves to polygons,
        and rasterises them onto a Pillow RGBA canvas.

        Args:
            svg_bytes: Raw SVG file content.

        Returns:
            PNG image bytes.
        """
        import io
        import xml.etree.ElementTree as ET

        from PIL import Image, ImageDraw
        from svgpathtools import parse_path

        root = ET.fromstring(svg_bytes)
        w = int(float(root.get("width", "100")))
        h = int(float(root.get("height", "100")))

        img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        draw = ImageDraw.Draw(img)

        paths = (
            root.findall("{http://www.w3.org/2000/svg}path")
            or root.findall("path")
        )

        for elem in paths:
            d = elem.get("d", "")
            fill = elem.get("fill", "none")
            transform = elem.get("transform", "")

            if fill == "none" or not d:
                continue
            if not (fill.startswith("#") and len(fill) >= 7):
                continue

            r_c = int(fill[1:3], 16)
            g_c = int(fill[3:5], 16)
            b_c = int(fill[5:7], 16)

            tx, ty = 0.0, 0.0
            tm = re.match(r"translate\(([^,]+),([^)]+)\)", transform)
            if tm:
                tx, ty = float(tm.group(1)), float(tm.group(2))

            try:
                path = parse_path(d)
                length = path.length()
                n = max(200, int(length))
                n = min(n, 5000)
                points = []
                for i in range(n):
                    pt = path.point(i / n)
                    points.append((pt.real + tx, pt.imag + ty))
                if len(points) >= 3:
                    draw.polygon(points, fill=(r_c, g_c, b_c, 255))
            except Exception:
                continue

        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()

    def _extract_and_save_content_images(
        self, html_str: str, assets_dir: str
    ) -> tuple[str, list[str]]:
        """Extract base64-embedded images from HTML, save as files, replace src with filename.

        This prevents huge base64 blobs from being sent to the LLM (where they
        would be truncated by ``_truncate_html``).  The saved files are later
        merged into the ``chart_images`` list so the LLM receives
        ``add_picture()`` instructions for them.

        Args:
            html_str: HTML that may contain ``<img src="data:image/...;base64,...">`` tags.
            assets_dir: Directory to write extracted image files into.

        Returns:
            Tuple of (cleaned_html, list_of_saved_filenames).
        """
        filenames: list[str] = []
        counter = 0

        def _replace(match: re.Match) -> str:
            nonlocal counter
            prefix = match.group(1)       # '<img ... src="'
            mime_sub = match.group(2)      # 'png', 'jpeg', etc.
            b64_data = match.group(3)      # raw base64 payload
            suffix = match.group(4)        # closing '"'

            ext = self._EXT_MAP.get(mime_sub.lower(), ".png")
            filename = f"content_image_{counter}{ext}"
            counter += 1

            # Decode and write
            try:
                image_bytes = base64.b64decode(b64_data)

                # python-pptx uses PIL which cannot handle SVG files.
                # Convert SVG to PNG so add_picture() works.
                if ext == ".svg":
                    image_bytes = self._svg_to_png(image_bytes)
                    filename = filename.rsplit(".", 1)[0] + ".png"

                filepath = Path(assets_dir) / filename
                filepath.write_bytes(image_bytes)
                filenames.append(filename)
                logger.info("Extracted content image", extra={"filename": filename, "size": len(image_bytes)})
            except Exception as e:
                logger.warning("Failed to extract content image", exc_info=True, extra={"error": str(e)})
                # Return original match unchanged on error
                return match.group(0)

            return f"{prefix}{filename}{suffix}"

        cleaned = self._BASE64_IMG_RE.sub(_replace, html_str)
        return cleaned, filenames

    def _truncate_html(self, html_str: str, max_length: int = 15000) -> str:
        """Truncate HTML to reasonable length for LLM while preserving CSS.
        
        Args:
            html_str: HTML string to truncate
            max_length: Maximum length before truncation
        
        Returns:
            Truncated HTML string with CSS preserved
        """
        if len(html_str) <= max_length:
            logger.debug(
                "HTML within length limit, no truncation needed",
                extra={"html_length": len(html_str), "max_length": max_length}
            )
            return html_str
        
        logger.debug(
            "Truncating HTML for LLM (preserving CSS)",
            extra={
                "original_length": len(html_str),
                "max_length": max_length,
                "html_preview_before": html_str[:500] + "..." if len(html_str) > 500 else html_str,
            }
        )
        
        try:
            soup = BeautifulSoup(html_str, 'lxml')
            
            # CRITICAL: Preserve style tag with all CSS for color extraction
            # Extract style content before truncation
            style_tags = soup.find_all('style')
            style_content = '\n'.join([str(tag) for tag in style_tags])
            
            # Remove script tags (charts will use screenshots or extracted data)
            script_count = len(soup.find_all('script'))
            for script in soup.find_all('script'):
                script.decompose()
            
            # Get body content
            body_content = str(soup.body) if soup.body else ""
            
            # If still too long, truncate body but keep style
            if len(body_content) + len(style_content) > max_length:
                # Keep style tag, truncate body content
                body_max = max_length - len(style_content) - 500  # Reserve space for style
                if body_max > 0:
                    body_content = body_content[:body_max] + "..."
            
            # Reconstruct with style tag preserved
            truncated = f"<html><head>{style_content}</head>{body_content}</html>"
            
            logger.debug(
                "HTML truncated successfully (CSS preserved)",
                extra={
                    "original_length": len(html_str),
                    "truncated_length": len(truncated),
                    "style_length": len(style_content),
                    "scripts_removed": script_count,
                    "truncated_preview": truncated[:500] + "..." if len(truncated) > 500 else truncated,
                }
            )
            
            return truncated
        except Exception as e:
            logger.warning(
                "HTML truncation fallback to simple slice",
                exc_info=True,
                extra={
                    "error": str(e),
                    "original_length": len(html_str),
                }
            )
            # Even in fallback, try to preserve style tag
            if '<style>' in html_str and '</style>' in html_str:
                style_start = html_str.find('<style>')
                style_end = html_str.find('</style>') + 8
                style_tag = html_str[style_start:style_end]
                body_slice = html_str[style_end:style_end + (max_length - len(style_tag) - 100)]
                return f"<html><head>{style_tag}</head><body>{body_slice}...</body></html>"
            return html_str[:max_length]
    
    @staticmethod
    def _sanitize_code(code: str) -> str:
        """Fix common LLM-generated code issues before execution.

        Handles smart quotes, em-dashes, and apostrophes inside string literals
        that cause SyntaxError (e.g. ``Anthony's`` inside a single-quoted string).
        """
        # Smart/curly quotes → straight quotes
        for smart, straight in {
            "\u2018": "'", "\u2019": "'",   # single curly quotes
            "\u201c": '"', "\u201d": '"',   # double curly quotes
        }.items():
            code = code.replace(smart, straight)

        # Em-dash / en-dash → regular hyphen (prevents SyntaxError when
        # they appear outside strings due to broken quoting)
        code = code.replace("\u2014", "-")   # em-dash
        code = code.replace("\u2013", "-")   # en-dash

        # Fix apostrophes inside single-quoted strings (e.g. 'Anthony's workflow')
        # by converting affected lines to use double-quoted strings.
        import ast as _ast
        for _ in range(20):
            try:
                _ast.parse(code)
                break
            except SyntaxError as exc:
                if exc.lineno is None:
                    break
                lines = code.splitlines()
                idx = exc.lineno - 1
                if idx < 0 or idx >= len(lines):
                    break
                original = lines[idx]
                # Re-quote single-quoted strings that contain apostrophes
                fixed = re.sub(
                    r"'([^']*\w'\w[^']*)'",
                    lambda m: '"' + m.group(1) + '"',
                    original,
                )
                if fixed == original:
                    break
                lines[idx] = fixed
                code = "\n".join(lines)
        return code

    def _execute_single_slide_converter(
        self,
        code: str,
        html_str: str,
        output_path: str,
        assets_dir: str
    ) -> None:
        """Execute generated converter code for single slide.
        
        Args:
            code: Generated Python code
            html_str: HTML content
            output_path: Path to save PPTX
            assets_dir: Directory containing assets
        
        Raises:
            PPTXConversionError: If code execution fails
        """
        code = self._sanitize_code(code)

        required_imports = """from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
"""
        if "from pptx import Presentation" not in code and "from pptx import" not in code:
            code = required_imports + "\n" + code
            logger.debug("Injected required imports into generated code")
        
        temp_module_path = Path(tempfile.mktemp(suffix=".py", prefix="converter_"))
        temp_module_path.write_text(code, encoding='utf-8')
        
        try:
            # Load as module
            spec = importlib.util.spec_from_file_location("temp_converter", str(temp_module_path))
            module = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(module)
            
            # Execute convert_to_pptx function
            module.convert_to_pptx(html_str, output_path, assets_dir)
            
        except Exception as e:
            logger.error("Failed to execute converter code", exc_info=True)
            raise PPTXConversionError(f"Code execution failed: {str(e)}") from e
        finally:
            # Cleanup temp module file
            if temp_module_path.exists():
                temp_module_path.unlink()
    
    def _execute_slide_adder(
        self,
        code: str,
        prs: Presentation,
        html_str: str,
        assets_dir: str
    ) -> None:
        """Execute generated code to add slide to presentation.
        
        Args:
            code: Generated Python code
            prs: Presentation object
            html_str: HTML content
            assets_dir: Directory containing assets
        """
        code = self._sanitize_code(code)

        required_imports = """from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
"""
        if "from pptx.util import" not in code:
            code = required_imports + "\n" + code
            logger.debug("Injected required imports into generated code")

        temp_module_path = Path(tempfile.mktemp(suffix=".py", prefix="slide_adder_"))
        temp_module_path.write_text(code, encoding='utf-8')

        # Also save to debug dir for inspection (use .txt to avoid triggering uvicorn reload)
        debug_dir = Path("logs/pptx_debug")
        debug_dir.mkdir(parents=True, exist_ok=True)
        debug_file = debug_dir / f"slide_{len(prs.slides) + 1}_code.txt"
        debug_file.write_text(code, encoding='utf-8')
        print(f"[PPTX_CONVERTER] Generated code saved to: {debug_file}")

        slides_before = len(prs.slides)

        try:
            # Load as module and execute the generated function
            spec = importlib.util.spec_from_file_location("temp_slide_adder", str(temp_module_path))
            module = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(module)
            module.add_slide_to_presentation(prs, html_str, assets_dir)

            if len(prs.slides) > slides_before:
                new_slide = prs.slides[len(prs.slides) - 1]
                pic_shapes = [s for s in new_slide.shapes if s.shape_type == 13]  # 13 = PICTURE
                print(f"[PPTX_CONVERTER] Slide has {len(new_slide.shapes)} shapes, {len(pic_shapes)} are pictures")
                if not pic_shapes and "add_picture" in code:
                    print(f"[PPTX_CONVERTER] WARNING: Code contains add_picture but no picture shapes found on slide!")
        except Exception as e:
            logger.warning(
                "Generated code error, creating fallback slide",
                exc_info=True,
                extra={"error": str(e)}
            )
            (debug_dir / f"slide_{slides_before + 1}_ERROR.txt").write_text(str(e), encoding='utf-8')

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = "Slide Content"
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.color.rgb = RGBColor(16, 32, 37)
        finally:
            # Cleanup temp module file
            if temp_module_path.exists():
                temp_module_path.unlink()
