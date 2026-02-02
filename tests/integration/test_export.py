"""Integration tests for export functionality (PPTX, HTML, PDF)."""

import io
import zipfile
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from fastapi.testclient import TestClient

from src.api.main import app


# =============================================================================
# Helper Functions
# =============================================================================


def create_mock_slide_deck(slide_count: int, with_charts: bool = False) -> dict:
    """Create a mock slide deck for testing.

    Args:
        slide_count: Number of slides to create
        with_charts: Whether to include chart canvases in slides

    Returns:
        Mock slide deck dictionary
    """
    slides = []
    for i in range(slide_count):
        slide = {
            "index": i,
            "slide_id": f"slide_{i}",
            "html": f"""<div class="slide">
                <h1>Slide {i + 1}</h1>
                <p>Content for slide {i + 1}</p>
            </div>""",
            "title": f"Slide {i + 1}",
            "verification_status": "verified",
        }

        if with_charts and i % 2 == 0:
            slide["html"] = f"""<div class="slide">
                <h1>Chart Slide {i + 1}</h1>
                <canvas id="chart_{i}"></canvas>
            </div>"""

        slides.append(slide)

    return {
        "slides": slides,
        "slide_count": slide_count,
        "title": "Test Presentation",
        "css": ".slide { width: 960px; height: 540px; }",
        "scripts": "// Chart initialization",
        "external_scripts": ["https://cdn.jsdelivr.net/npm/chart.js"],
    }


# =============================================================================
# Fixtures
# =============================================================================


@pytest.fixture
def client():
    """Create test client."""
    with TestClient(app) as test_client:
        yield test_client


@pytest.fixture
def mock_chat_service():
    """Mock the chat service."""
    with patch("src.api.routes.export.get_chat_service") as mock:
        service = MagicMock()
        mock.return_value = service
        yield service


@pytest.fixture
def mock_pptx_converter():
    """Mock the PPTX converter."""
    with patch("src.api.routes.export.HtmlToPptxConverterV3") as mock_class:
        converter = MagicMock()
        # Make convert_slide_deck an async mock
        converter.convert_slide_deck = AsyncMock()
        mock_class.return_value = converter
        yield converter


# =============================================================================
# PPTX Export Tests
# =============================================================================


class TestPPTXExport:
    """Tests for PowerPoint export."""

    def test_pptx_export_requires_session_id(self, client):
        """POST /api/export/pptx returns 422 without session_id."""
        response = client.post("/api/export/pptx", json={})
        assert response.status_code == 422

    def test_pptx_export_session_not_found(self, client, mock_chat_service):
        """POST /api/export/pptx returns 404 for missing session."""
        mock_chat_service.get_slides.return_value = None

        response = client.post("/api/export/pptx", json={"session_id": "nonexistent"})
        assert response.status_code == 404

    def test_pptx_export_no_slides(self, client, mock_chat_service):
        """POST /api/export/pptx returns 404 when no slides exist (empty deck)."""
        mock_chat_service.get_slides.return_value = {"slides": [], "slide_count": 0}

        response = client.post("/api/export/pptx", json={"session_id": "test-123"})
        # Based on the actual implementation, empty slides returns 404 "No slides available"
        assert response.status_code == 404
        assert "no slides" in response.json()["detail"].lower()

    def test_pptx_export_success(
        self, client, mock_chat_service, mock_pptx_converter, tmp_path
    ):
        """POST /api/export/pptx returns PPTX file."""
        mock_chat_service.get_slides.return_value = create_mock_slide_deck(3)

        # Create a valid PPTX file for the mock to "create"
        # PPTX files are ZIP archives with specific structure
        from pptx import Presentation

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(pptx_path))

        # Make the converter write to the output path
        async def mock_convert(
            slides, output_path, use_screenshot=True, html_source_paths=None, chart_images_per_slide=None
        ):
            # Copy our test PPTX to the output path
            import shutil
            shutil.copy(str(pptx_path), output_path)

        mock_pptx_converter.convert_slide_deck.side_effect = mock_convert

        response = client.post("/api/export/pptx", json={"session_id": "test-123"})

        assert response.status_code == 200
        assert (
            response.headers["content-type"]
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        assert "content-disposition" in response.headers

    def test_pptx_export_filename(
        self, client, mock_chat_service, mock_pptx_converter, tmp_path
    ):
        """PPTX export has correct filename in Content-Disposition."""
        deck = create_mock_slide_deck(3)
        deck["title"] = "My Test Presentation"
        mock_chat_service.get_slides.return_value = deck

        # Create a valid PPTX file
        from pptx import Presentation

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(pptx_path))

        async def mock_convert(
            slides, output_path, use_screenshot=True, html_source_paths=None, chart_images_per_slide=None
        ):
            import shutil
            shutil.copy(str(pptx_path), output_path)

        mock_pptx_converter.convert_slide_deck.side_effect = mock_convert

        response = client.post("/api/export/pptx", json={"session_id": "test-123"})

        disposition = response.headers.get("content-disposition", "")
        assert "attachment" in disposition
        assert ".pptx" in disposition

    def test_pptx_export_with_chart_images(
        self, client, mock_chat_service, mock_pptx_converter, tmp_path
    ):
        """PPTX export accepts chart images for chart slides."""
        mock_chat_service.get_slides.return_value = create_mock_slide_deck(
            2, with_charts=True
        )

        # Create a valid PPTX file
        from pptx import Presentation

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(pptx_path))

        async def mock_convert(
            slides, output_path, use_screenshot=True, html_source_paths=None, chart_images_per_slide=None
        ):
            import shutil
            shutil.copy(str(pptx_path), output_path)

        mock_pptx_converter.convert_slide_deck.side_effect = mock_convert

        response = client.post(
            "/api/export/pptx",
            json={
                "session_id": "test-123",
                "use_screenshot": True,
                "chart_images": [
                    [
                        {
                            "canvas_id": "chart_0",
                            "base64_data": "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg==",
                        }
                    ],
                    [],
                ],
            },
        )

        assert response.status_code == 200

    def test_pptx_export_conversion_error(
        self, client, mock_chat_service, mock_pptx_converter
    ):
        """PPTX export handles conversion errors gracefully."""
        from src.services.html_to_pptx import PPTXConversionError

        mock_chat_service.get_slides.return_value = create_mock_slide_deck(3)
        mock_pptx_converter.convert_slide_deck.side_effect = PPTXConversionError(
            "Font not found"
        )

        response = client.post("/api/export/pptx", json={"session_id": "test-123"})

        assert response.status_code == 500
        assert "conversion" in response.json()["detail"].lower()

    def test_pptx_content_is_valid_zip(
        self, client, mock_chat_service, mock_pptx_converter, tmp_path
    ):
        """PPTX export produces valid ZIP file (PPTX is ZIP-based)."""
        mock_chat_service.get_slides.return_value = create_mock_slide_deck(2)

        # Create a valid PPTX file
        from pptx import Presentation

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(pptx_path))

        async def mock_convert(
            slides, output_path, use_screenshot=True, html_source_paths=None, chart_images_per_slide=None
        ):
            import shutil
            shutil.copy(str(pptx_path), output_path)

        mock_pptx_converter.convert_slide_deck.side_effect = mock_convert

        response = client.post("/api/export/pptx", json={"session_id": "test-123"})

        if response.status_code == 200:
            # PPTX files are ZIP archives
            content = io.BytesIO(response.content)
            assert zipfile.is_zipfile(content)

    def test_pptx_export_with_use_screenshot_false(
        self, client, mock_chat_service, mock_pptx_converter, tmp_path
    ):
        """PPTX export respects use_screenshot=False."""
        mock_chat_service.get_slides.return_value = create_mock_slide_deck(2)

        # Create a valid PPTX file
        from pptx import Presentation

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(pptx_path))

        captured_args = {}

        async def mock_convert(
            slides, output_path, use_screenshot=True, html_source_paths=None, chart_images_per_slide=None
        ):
            import shutil
            captured_args["use_screenshot"] = use_screenshot
            shutil.copy(str(pptx_path), output_path)

        mock_pptx_converter.convert_slide_deck.side_effect = mock_convert

        response = client.post(
            "/api/export/pptx",
            json={"session_id": "test-123", "use_screenshot": False},
        )

        assert response.status_code == 200
        assert captured_args.get("use_screenshot") is False


# =============================================================================
# HTML Export Tests (if endpoint exists)
# =============================================================================


class TestHTMLExport:
    """Tests for HTML export (if implemented).

    Note: The current export.py may not have an HTML export endpoint.
    These tests verify behavior if/when HTML export is added.
    """

    def test_html_export_endpoint_existence(self, client):
        """Check if HTML export endpoint exists."""
        # Try to access HTML export endpoint
        response = client.get("/api/export/html?session_id=test-123")

        # If endpoint doesn't exist, it returns 404 or 405
        # If it exists but session not found, it would return a different error
        # This test documents whether the endpoint is implemented
        if response.status_code == 404:
            pytest.skip("HTML export endpoint not implemented")
        elif response.status_code == 405:
            pytest.skip("HTML export endpoint not implemented (method not allowed)")


# =============================================================================
# PDF Export Tests
# =============================================================================


class TestPDFExport:
    """Tests for PDF export (if implemented)."""

    def test_pdf_export_endpoint_existence(self, client):
        """Check if PDF export endpoint exists."""
        response = client.post("/api/export/pdf", json={"session_id": "test-123"})

        # PDF export may not be implemented
        if response.status_code in [404, 405]:
            pytest.skip("PDF export endpoint not implemented")
        elif response.status_code == 422:
            # Endpoint exists but validation failed - this is expected
            pass


# =============================================================================
# Export Error Handling Tests
# =============================================================================


class TestExportErrors:
    """Tests for export error handling."""

    def test_export_with_invalid_session(self, client, mock_chat_service):
        """Export returns 404 for invalid session."""
        mock_chat_service.get_slides.return_value = None

        response = client.post(
            "/api/export/pptx", json={"session_id": "invalid-session"}
        )
        assert response.status_code == 404

    def test_export_with_empty_deck(self, client, mock_chat_service):
        """Export returns 404 for empty deck."""
        mock_chat_service.get_slides.return_value = {"slides": [], "slide_count": 0}

        response = client.post("/api/export/pptx", json={"session_id": "test-123"})
        # The actual implementation returns 404 for no slides
        assert response.status_code == 404

    def test_export_with_malformed_html(
        self, client, mock_chat_service, mock_pptx_converter, tmp_path
    ):
        """Export handles malformed slide HTML gracefully."""
        deck = create_mock_slide_deck(1)
        deck["slides"][0]["html"] = "<div>Unclosed"
        mock_chat_service.get_slides.return_value = deck

        # Create a valid PPTX file
        from pptx import Presentation

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(pptx_path))

        async def mock_convert(
            slides, output_path, use_screenshot=True, html_source_paths=None, chart_images_per_slide=None
        ):
            import shutil
            shutil.copy(str(pptx_path), output_path)

        mock_pptx_converter.convert_slide_deck.side_effect = mock_convert

        # Should either succeed with best-effort or return error
        response = client.post("/api/export/pptx", json={"session_id": "test-123"})
        assert response.status_code in [200, 400, 500]

    def test_export_general_exception(self, client, mock_chat_service, mock_pptx_converter):
        """Export handles general exceptions gracefully."""
        mock_chat_service.get_slides.return_value = create_mock_slide_deck(3)
        mock_pptx_converter.convert_slide_deck.side_effect = Exception(
            "Unexpected error"
        )

        response = client.post("/api/export/pptx", json={"session_id": "test-123"})
        assert response.status_code == 500
        assert "failed" in response.json()["detail"].lower()


# =============================================================================
# Large Deck Export Tests
# =============================================================================


class TestLargeDeckExport:
    """Tests for exporting large decks."""

    def test_export_large_deck(
        self, client, mock_chat_service, mock_pptx_converter, tmp_path
    ):
        """Export works for large decks (20+ slides)."""
        mock_chat_service.get_slides.return_value = create_mock_slide_deck(25)

        # Create a valid PPTX file
        from pptx import Presentation

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(pptx_path))

        async def mock_convert(
            slides, output_path, use_screenshot=True, html_source_paths=None, chart_images_per_slide=None
        ):
            import shutil
            shutil.copy(str(pptx_path), output_path)

        mock_pptx_converter.convert_slide_deck.side_effect = mock_convert

        response = client.post("/api/export/pptx", json={"session_id": "test-123"})
        assert response.status_code == 200

    def test_export_deck_with_many_charts(
        self, client, mock_chat_service, mock_pptx_converter, tmp_path
    ):
        """Export works for decks with many charts."""
        mock_chat_service.get_slides.return_value = create_mock_slide_deck(
            10, with_charts=True
        )

        # Create a valid PPTX file
        from pptx import Presentation

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(pptx_path))

        async def mock_convert(
            slides, output_path, use_screenshot=True, html_source_paths=None, chart_images_per_slide=None
        ):
            import shutil
            shutil.copy(str(pptx_path), output_path)

        mock_pptx_converter.convert_slide_deck.side_effect = mock_convert

        response = client.post(
            "/api/export/pptx",
            json={"session_id": "test-123", "use_screenshot": True},
        )
        assert response.status_code == 200


# =============================================================================
# Async Export Tests
# =============================================================================


class TestAsyncPPTXExport:
    """Tests for async PPTX export endpoints."""

    def test_async_export_requires_session_id(self, client):
        """POST /api/export/pptx/async returns 422 without session_id."""
        response = client.post("/api/export/pptx/async", json={})
        assert response.status_code == 422

    def test_async_export_session_not_found(self, client, mock_chat_service):
        """POST /api/export/pptx/async returns 404 for missing session."""
        mock_chat_service.get_slides.return_value = None

        response = client.post(
            "/api/export/pptx/async", json={"session_id": "nonexistent"}
        )
        assert response.status_code == 404

    def test_async_export_no_slides(self, client, mock_chat_service):
        """POST /api/export/pptx/async returns 404 when no slides exist."""
        mock_chat_service.get_slides.return_value = {"slides": [], "slide_count": 0}

        response = client.post(
            "/api/export/pptx/async", json={"session_id": "test-123"}
        )
        assert response.status_code == 404

    def test_async_export_starts_job(self, client, mock_chat_service):
        """POST /api/export/pptx/async starts an export job."""
        mock_chat_service.get_slides.return_value = create_mock_slide_deck(3)

        with patch("src.api.services.export_job_queue.enqueue_export_job") as mock_enqueue:
            with patch("src.api.services.export_job_queue.generate_job_id", return_value="test-job-123"):
                mock_enqueue.return_value = None

                response = client.post(
                    "/api/export/pptx/async", json={"session_id": "test-123"}
                )

                assert response.status_code == 200
                data = response.json()
                assert data["job_id"] == "test-job-123"
                assert data["status"] == "pending"
                assert data["total_slides"] == 3

    def test_poll_export_job_not_found(self, client):
        """GET /api/export/pptx/poll/{job_id} returns 404 for missing job."""
        with patch("src.api.services.export_job_queue.get_export_job_status", return_value=None):
            response = client.get("/api/export/pptx/poll/nonexistent-job")
            assert response.status_code == 404

    def test_poll_export_job_in_progress(self, client):
        """GET /api/export/pptx/poll/{job_id} returns progress for running job."""
        job_status = {
            "status": "running",
            "progress": 5,
            "total_slides": 10,
            "error": None,
        }

        with patch(
            "src.api.services.export_job_queue.get_export_job_status", return_value=job_status
        ):
            response = client.get("/api/export/pptx/poll/test-job-123")

            assert response.status_code == 200
            data = response.json()
            assert data["status"] == "running"
            assert data["progress"] == 5
            assert data["total_slides"] == 10

    def test_download_export_job_not_found(self, client):
        """GET /api/export/pptx/download/{job_id} returns 404 for missing job."""
        with patch("src.api.services.export_job_queue.get_export_job_status", return_value=None):
            response = client.get("/api/export/pptx/download/nonexistent-job")
            assert response.status_code == 404

    def test_download_export_job_not_ready(self, client):
        """GET /api/export/pptx/download/{job_id} returns 400 if not completed."""
        job_status = {
            "status": "running",
            "progress": 5,
            "total_slides": 10,
            "error": None,
        }

        with patch(
            "src.api.services.export_job_queue.get_export_job_status", return_value=job_status
        ):
            response = client.get("/api/export/pptx/download/test-job-123")

            assert response.status_code == 400
            assert "not ready" in response.json()["detail"].lower()


# =============================================================================
# Build Slide HTML Function Tests
# =============================================================================


class TestBuildSlideHTML:
    """Tests for the build_slide_html helper function."""

    def test_build_slide_html_basic(self):
        """Test building HTML for a basic slide."""
        from src.api.routes.export import build_slide_html

        slide = {
            "slide_id": "test-slide",
            "html": "<div>Test content</div>",
        }
        slide_deck = {
            "title": "Test Deck",
            "css": ".test { color: red; }",
            "scripts": "",
            "external_scripts": [],
        }

        html = build_slide_html(slide, slide_deck)

        assert "<!DOCTYPE html>" in html
        assert "<html" in html
        assert "</html>" in html
        assert "Test content" in html
        assert ".test { color: red; }" in html

    def test_build_slide_html_with_external_scripts(self):
        """Test building HTML with external scripts."""
        from src.api.routes.export import build_slide_html

        slide = {
            "slide_id": "test-slide",
            "html": "<canvas id='chart'></canvas>",
        }
        slide_deck = {
            "title": "Test Deck",
            "css": "",
            "scripts": "",
            "external_scripts": [
                "https://cdn.jsdelivr.net/npm/chart.js",
                "https://example.com/other.js",
            ],
        }

        html = build_slide_html(slide, slide_deck)

        assert "cdn.jsdelivr.net/npm/chart.js" in html
        assert "example.com/other.js" in html
        assert "<script src=" in html

    def test_build_slide_html_with_scripts(self):
        """Test building HTML with inline scripts."""
        from src.api.routes.export import build_slide_html

        slide = {
            "slide_id": "test-slide",
            "html": "<canvas id='chart'></canvas>",
        }
        slide_deck = {
            "title": "Test Deck",
            "css": "",
            "scripts": "console.log('chart init');",
            "external_scripts": [],
        }

        html = build_slide_html(slide, slide_deck)

        assert "console.log('chart init');" in html

    def test_build_slide_html_preserves_slide_content(self):
        """Test that slide HTML content is preserved."""
        from src.api.routes.export import build_slide_html

        slide = {
            "slide_id": "test-slide",
            "html": '<div class="slide"><h1>Title</h1><p>Content</p></div>',
        }
        slide_deck = {
            "title": "Test",
            "css": "",
            "scripts": "",
            "external_scripts": [],
        }

        html = build_slide_html(slide, slide_deck)

        assert '<div class="slide">' in html
        assert "<h1>Title</h1>" in html
        assert "<p>Content</p>" in html
