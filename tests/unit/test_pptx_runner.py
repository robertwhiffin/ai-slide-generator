# tests/unit/test_pptx_runner.py
"""Tests for the in-jail PPTX runner (SDR-4437 PR-5).

Most tests call the runner's importable functions directly (no subprocess) so
the per-slide loop, fallback isolation, and progress emission are unit-testable
on any platform. NOTE: the direct-call path passes an explicit `emit=` callback,
which SKIPS run()'s fd-level stdout redirection (that block only runs when
emit is None) — so the anti-spoofing control itself is exercised by the
subprocess-level test at the bottom, which goes through the real jail.
"""

import json
from pathlib import Path

from pptx import Presentation

from src.services.converter_jail import pptx_runner, protocol


def _write_slide(job_dir: Path, index: int, code, html: str = "<p>x</p>"):
    sdir = job_dir / f"slide_{index:03d}"
    (sdir / protocol.ASSETS_DIR).mkdir(parents=True)
    (sdir / protocol.HTML_NAME).write_text(html)
    if code is not None:
        (sdir / protocol.CODE_NAME).write_text(code)
    return {"index": index, "has_code": code is not None, "dir": sdir.name}


_GOOD = (
    "def add_slide_to_presentation(prs, html_str, assets_dir):\n"
    "    from pptx.util import Inches\n"
    "    slide = prs.slides.add_slide(prs.slide_layouts[6])\n"
    "    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))\n"
    "    box.text_frame.text = 'ok'\n"
)
_BAD = (
    "def add_slide_to_presentation(prs, html_str, assets_dir):\n"
    "    raise RuntimeError('boom')\n"
)


def test_builds_all_slides_with_fallback_isolation(tmp_path):
    job = tmp_path / "job"
    job.mkdir()
    manifest = {"slides": [
        _write_slide(job, 0, _GOOD),
        _write_slide(job, 1, _BAD),      # raises -> fallback, deck survives
        _write_slide(job, 2, None),      # no code -> fallback
    ]}
    (job / protocol.MANIFEST_NAME).write_text(json.dumps(manifest))
    out = tmp_path / "deck.pptx"

    events = []
    pptx_runner.run(str(job), str(out), emit=lambda line: events.append(line))

    prs = Presentation(str(out))
    assert len(prs.slides) == 3  # one bad slide degraded, not the whole deck

    progress = [protocol.decode_progress(e) for e in events]
    progress = [p for p in progress if p]
    assert progress[-1][0] == 3 and progress[-1][1] == 3  # final (3, 3, ...)


def test_generated_print_does_not_reach_progress_channel(tmp_path):
    # NOTE: this covers only the emit-callback plumbing (generated code has no
    # handle on `emit`). It does NOT exercise the fd-level stdout redirection —
    # run() skips that block when emit is not None. The redirection itself is
    # proven by test_spoofed_progress_blocked_in_real_jail below.
    job = tmp_path / "job"
    job.mkdir()
    noisy = (
        "def add_slide_to_presentation(prs, html_str, assets_dir):\n"
        "    print('" + protocol.PROGRESS_PREFIX + "{\\\"current\\\": 99}')\n"
        "    slide = prs.slides.add_slide(prs.slide_layouts[6])\n"
    )
    manifest = {"slides": [_write_slide(job, 0, noisy)]}
    (job / protocol.MANIFEST_NAME).write_text(json.dumps(manifest))
    out = tmp_path / "deck.pptx"

    events = []
    pptx_runner.run(str(job), str(out), emit=lambda line: events.append(line))
    # Only the runner's own single (1,1,...) progress line should appear.
    decoded = [protocol.decode_progress(e) for e in events]
    decoded = [d for d in decoded if d]
    assert all(d[0] == 1 and d[1] == 1 for d in decoded)


def test_spoofed_progress_blocked_in_real_jail(tmp_path):
    """Subprocess-level anti-spoofing: generated code prints a FULLY
    well-formed progress line; because run() (emit=None path) redirects
    Python-level stdout to stderr via os.dup2 before any snippet executes,
    the spoof must never reach the host's progress_cb — only the runner's
    own line may."""
    from src.services.converter_jail import run_pptx_jail

    spoof = protocol.encode_progress(99, 99, "spoofed").rstrip("\n")
    noisy = (
        "def add_slide_to_presentation(prs, html_str, assets_dir):\n"
        f"    print({spoof!r})\n"
        "    slide = prs.slides.add_slide(prs.slide_layouts[6])\n"
    )
    job = tmp_path / "job"
    job.mkdir()
    manifest = {"slides": [_write_slide(job, 0, noisy)]}
    (job / protocol.MANIFEST_NAME).write_text(json.dumps(manifest))
    out = tmp_path / "deck.pptx"

    events = []
    result = run_pptx_jail(
        str(job), str(out),
        progress_cb=lambda c, t, m: events.append((c, t, m)),
    )

    assert result.timed_out is False
    assert result.returncode == 0, result.stderr_tail
    assert all(e[:2] == (1, 1) for e in events)   # no (99, 99, "spoofed")
    assert (1, 1, "Building slide 1/1…") in events  # real line got through
    assert len(Presentation(str(out)).slides) == 1


def test_runner_reports_which_slides_actually_rendered(tmp_path):
    """The trusted per-slide RESULT channel.

    Both fallback branches save a VALID slide, so the finished .pptx cannot tell
    the host whether anything real was produced. The runner therefore reports each
    slide's outcome, and this is what lets an all-placeholder deck be refused
    instead of returned as a successful export.
    """
    job = tmp_path / "job"
    job.mkdir()
    manifest = {"slides": [
        _write_slide(job, 0, _GOOD),   # renders
        _write_slide(job, 1, _BAD),    # raises -> placeholder
        _write_slide(job, 2, None),    # no code -> placeholder
    ]}
    (job / protocol.MANIFEST_NAME).write_text(json.dumps(manifest))

    events = []
    pptx_runner.run(str(job), str(tmp_path / "deck.pptx"), emit=lambda line: events.append(line))

    results = [protocol.decode_slide_result(e) for e in events]
    assert [r for r in results if r] == [(0, True), (1, False), (2, False)]


def test_slide_results_do_not_disturb_the_progress_channel(tmp_path):
    """Two sentinels share one stdout stream; each must decode as only itself."""
    job = tmp_path / "job"
    job.mkdir()
    (job / protocol.MANIFEST_NAME).write_text(
        json.dumps({"slides": [_write_slide(job, 0, _GOOD)]})
    )

    events = []
    pptx_runner.run(str(job), str(tmp_path / "deck.pptx"), emit=lambda line: events.append(line))

    progress = [p for p in (protocol.decode_progress(e) for e in events) if p]
    results = [r for r in (protocol.decode_slide_result(e) for e in events) if r]
    assert progress == [(1, 1, "Building slide 1/1…")]
    assert results == [(0, True)]
    # A result line must never decode as progress, or a slide report would be
    # relayed to the user as a progress event.
    assert protocol.decode_progress(protocol.encode_slide_result(0, True)) is None
    assert protocol.decode_slide_result(protocol.encode_progress(1, 1, "x")) is None


def test_the_host_collects_slide_results_through_the_real_jail(tmp_path):
    """End-to-end over the subprocess boundary: the outcomes reach JailResult."""
    from src.services.converter_jail import run_pptx_jail

    job = tmp_path / "job"
    job.mkdir()
    manifest = {"slides": [_write_slide(job, 0, _GOOD), _write_slide(job, 1, _BAD)]}
    (job / protocol.MANIFEST_NAME).write_text(json.dumps(manifest))

    result = run_pptx_jail(str(job), str(tmp_path / "deck.pptx"))

    assert result.returncode == 0, result.stderr_tail
    assert result.slide_results == (True, False)


def test_a_snippet_cannot_forge_its_own_slide_result(tmp_path):
    """Same anti-spoofing property the progress channel has, for the same reason:
    a snippet that printed a success report could hide a content-free deck."""
    from src.services.converter_jail import run_pptx_jail

    spoof = protocol.encode_slide_result(0, True).rstrip("\n")
    forging = (
        "def add_slide_to_presentation(prs, html_str, assets_dir):\n"
        f"    print({spoof!r})\n"
        "    raise RuntimeError('no real content')\n"
    )
    job = tmp_path / "job"
    job.mkdir()
    (job / protocol.MANIFEST_NAME).write_text(
        json.dumps({"slides": [_write_slide(job, 0, forging)]})
    )

    result = run_pptx_jail(str(job), str(tmp_path / "deck.pptx"))

    assert result.returncode == 0, result.stderr_tail
    # Exactly ONE report, the runner's own, and it says the slide did NOT render.
    assert result.slide_results == (False,)
